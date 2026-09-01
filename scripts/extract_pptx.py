"""Pure-text extraction from a .pptx file: no LLM involved, no hallucination risk."""
import warnings

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _flatten_shapes(shapes):
    """Recursively expand group shapes so nested text boxes/tables are visible."""
    flat = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            flat.extend(_flatten_shapes(shape.shapes))
        else:
            flat.append(shape)
    return flat


def _shape_text(shape):
    lines = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if not text:
                text = para.text.strip()
            if text:
                lines.append(text)
    elif shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))
    return lines


def extract_slides(pptx_path):
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_lines = []

        shapes = _flatten_shapes(slide.shapes)
        # Reading order, not insertion/z-order: top-to-bottom, left-to-right.
        shapes.sort(key=lambda s: (s.top if s.top is not None else 0,
                                    s.left if s.left is not None else 0))

        for shape in shapes:
            if not (shape.has_text_frame or shape.has_table):
                continue
            try:
                is_title = shape == slide.shapes.title
                lines = _shape_text(shape)
            except Exception as e:
                warnings.warn(f"{pptx_path}: slide {i}: failed to read shape: {e}")
                continue
            if is_title and lines:
                title = lines[0]
                body_lines.extend(lines[1:])
            else:
                body_lines.extend(lines)

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        if not title and not body_lines and not notes:
            continue

        slides.append({
            "slide_number": i,
            "title": title,
            "content": body_lines,
            "notes": notes,
        })
    return slides


if __name__ == "__main__":
    import sys
    import json
    slides = extract_slides(sys.argv[1])
    print(json.dumps(slides, indent=2, ensure_ascii=False))
